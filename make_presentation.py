from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

DARK   = RGBColor(0x1a, 0x1a, 0x1a)
WHITE  = RGBColor(0xf2, 0xf0, 0xec)
BLUE   = RGBColor(0x29, 0x80, 0xb9)
GREEN  = RGBColor(0x27, 0xae, 0x60)
ORANGE = RGBColor(0xe6, 0x7e, 0x22)
RED    = RGBColor(0xe7, 0x4c, 0x3c)
GRAY   = RGBColor(0x88, 0x88, 0x88)
LIGHT  = RGBColor(0xf5, 0xf4, 0xf1)
PURPLE = RGBColor(0x8e, 0x44, 0xad)

blank = prs.slide_layouts[6]  # completely blank

def add_rect(slide, x, y, w, h, fill=None, line=None):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, x, y, w, h, size=14, bold=False, color=DARK, align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Segoe UI"
    return txb

def label(slide, text, x, y, w=4):
    add_text(slide, text.upper(), x, y, w, 0.25, size=8, bold=True, color=GRAY)

def tag(slide, text, x, y, fill, tcolor=WHITE):
    r = add_rect(slide, x, y, len(text)*0.085+0.2, 0.28, fill=fill)
    add_text(slide, text, x+0.06, y+0.04, len(text)*0.085+0.1, 0.22, size=8, bold=True, color=tcolor, align=PP_ALIGN.CENTER)

# ─── SLIDE 1 — Title ──────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
add_rect(s, 0, 0, 13.33, 7.5, fill=DARK)
add_rect(s, 0, 0, 0.06, 7.5, fill=BLUE)

add_text(s, "ТЗ TOOL", 0.6, 1.4, 10, 1.4, size=72, bold=True, color=WHITE)
add_text(s, "Перетворює пакет файлів від клієнта\nна структуроване технічне завдання для 3D-команди",
         0.6, 3.1, 9, 1.2, size=20, color=RGBColor(0x66,0x66,0x66))

tag(s, "Web", 0.6, 4.6, BLUE)
tag(s, "macOS", 1.5, 4.6, GREEN)
tag(s, "Windows", 2.45, 4.6, GREEN)
tag(s, "Claude Sonnet 4.6", 3.55, 4.6, RGBColor(0x33,0x33,0x33))

add_text(s, "dlytvynov-rgb.github.io/tz-tool", 0.6, 6.4, 8, 0.4, size=11, color=GRAY)

# ─── SLIDE 2 — Problem ────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
add_rect(s, 0, 0, 13.33, 0.8, fill=DARK)
add_text(s, "ПРОБЛЕМА", 0.5, 0.2, 6, 0.5, size=11, bold=True, color=WHITE)

add_text(s, "Клієнт надсилає пакет файлів", 0.5, 1.2, 12, 0.6, size=28, bold=True, color=DARK)

files = [("📄", "Бриф\nу Word"), ("📐", "Креслення\nDWG"), ("🖼", "Референси\nPDF"), ("📊", "Таблиці\nExcel"), ("💬", "Фото і\nнотатки")]
for i, (ico, lbl) in enumerate(files):
    x = 0.5 + i * 2.4
    add_rect(s, x, 2.1, 2.1, 1.6, fill=LIGHT, line=RGBColor(0xe0,0xe0,0xe0))
    add_text(s, ico, x+0.7, 2.2, 0.8, 0.6, size=28, align=PP_ALIGN.CENTER)
    add_text(s, lbl, x+0.1, 2.9, 1.9, 0.7, size=11, color=GRAY, align=PP_ALIGN.CENTER)

add_rect(s, 0.5, 4.1, 12.3, 0.06, fill=RGBColor(0xee,0xee,0xee))

add_text(s, "Вручну читати, зіставляти і складати ТЗ —", 0.5, 4.4, 8, 0.5, size=18, color=DARK)
add_text(s, "1–3 години на один проект", 0.5, 4.95, 8, 0.5, size=22, bold=True, color=RED)

add_rect(s, 8.8, 4.2, 4, 1.6, fill=RGBColor(0xff,0xf3,0xf3), line=RED)
add_text(s, "Вимоги губляться\nСуперечності не помічають\nПереробки в середині роботи", 9.0, 4.35, 3.6, 1.3, size=12, color=RED)

# ─── SLIDE 3 — Solution ───────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
add_rect(s, 0, 0, 13.33, 0.8, fill=DARK)
add_text(s, "РІШЕННЯ", 0.5, 0.2, 6, 0.5, size=11, bold=True, color=WHITE)

add_text(s, "Одна кнопка — повне ТЗ", 0.5, 1.1, 12, 0.7, size=32, bold=True, color=DARK)

steps = [
    (BLUE,   "01", "Завантаж файли",     "PDF, DWG, DOCX, Excel,\nзображення, ZIP-архіви"),
    (BLUE,   "02", "Натисни кнопку",     "\"Розібрати файли\" —\nодин Claude-запит"),
    (GREEN,  "03", "Отримай результат",  "Структуроване ТЗ готове\nдо передачі в роботу"),
]
for i, (clr, num, title, desc) in enumerate(steps):
    x = 0.5 + i * 4.1
    add_rect(s, x, 2.1, 3.7, 2.8, fill=LIGHT, line=RGBColor(0xe0,0xe0,0xe0))
    add_rect(s, x, 2.1, 3.7, 0.45, fill=clr)
    add_text(s, num, x+0.15, 2.15, 0.5, 0.35, size=14, bold=True, color=WHITE)
    add_text(s, title, x+0.6, 2.15, 3.0, 0.35, size=13, bold=True, color=WHITE)
    add_text(s, desc, x+0.2, 2.75, 3.3, 1.0, size=12, color=DARK)

add_rect(s, 0.5, 5.3, 12.3, 1.6, fill=RGBColor(0xf0,0xf8,0xf0), line=GREEN)
add_text(s, "✓  Вимоги по кімнатах і стадіях виробництва\n✓  Виявлені суперечності між файлами\n✓  Список питань до клієнта з дефолтами\n✓  Роадмап і чеклист здачі",
         0.8, 5.4, 11.5, 1.4, size=13, color=RGBColor(0x1a,0x5c,0x1a))

# ─── SLIDE 4 — Input formats ──────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
add_rect(s, 0, 0, 13.33, 0.8, fill=DARK)
add_text(s, "ВХІДНІ ФАЙЛИ", 0.5, 0.2, 6, 0.5, size=11, bold=True, color=WHITE)
add_text(s, "Підтримувані формати", 0.5, 1.1, 12, 0.6, size=28, bold=True, color=DARK)

fmts = [
    ("📄", "PDF", ".pdf — брифи, презентації,\nреференси"),
    ("📝", "Word", ".docx — текстові ТЗ,\nбрифи клієнта"),
    ("📐", "DWG / DXF", ".dwg .dxf — креслення,\nплани приміщень"),
    ("📊", "Excel", ".xlsx .csv — специфікації,\nтаблиці матеріалів"),
    ("🖼", "Зображення", ".jpg .png .webp —\nреференси, фото"),
    ("🗜", "ZIP-архів", ".zip — пакет файлів,\nрозпаковується авто"),
]
for i, (ico, name, desc) in enumerate(fmts):
    col, row = i % 3, i // 3
    x = 0.5 + col * 4.1
    y = 2.0 + row * 2.3
    add_rect(s, x, y, 3.8, 2.0, fill=LIGHT, line=RGBColor(0xe0,0xe0,0xe0))
    add_text(s, ico, x+0.2, y+0.2, 0.7, 0.6, size=24)
    add_text(s, name, x+0.95, y+0.25, 2.7, 0.4, size=14, bold=True, color=DARK)
    add_text(s, desc, x+0.2, y+0.85, 3.4, 0.9, size=11, color=GRAY)

# ─── SLIDE 5 — Views ──────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
add_rect(s, 0, 0, 13.33, 0.8, fill=DARK)
add_text(s, "РЕЖИМИ ПЕРЕГЛЯДУ", 0.5, 0.2, 8, 0.5, size=11, bold=True, color=WHITE)
add_text(s, "Чотири способи дивитись на результат", 0.5, 1.1, 12, 0.6, size=26, bold=True, color=DARK)

views = [
    (DARK,   "КІМНАТИ",  "Вимоги згруповані\nпо приміщеннях.\nОсновний режим."),
    (DARK,   "СТАДІЇ",   "Ті ж вимоги але\nпо виробничих стадіях:\nМоделінг → Видача."),
    (BLUE,   "ТАБЛИЦЯ",  "Зведена таблиця\nз фільтрами і сортуванням.\nЕкспорт в XLS."),
    (PURPLE, "ЗВІТ",     "SOW-матриця:\n✅ знайдено / ⚠️ неповно\n/ ❌ відсутнє."),
]
for i, (clr, name, desc) in enumerate(views):
    x = 0.5 + i * 3.1
    add_rect(s, x, 2.0, 2.8, 4.8, fill=LIGHT, line=RGBColor(0xe0,0xe0,0xe0))
    add_rect(s, x, 2.0, 2.8, 0.45, fill=clr)
    add_text(s, name, x+0.15, 2.1, 2.5, 0.3, size=11, bold=True, color=WHITE)
    add_text(s, desc, x+0.2, 2.65, 2.4, 3.8, size=12, color=DARK)

# ─── SLIDE 6 — Sections ───────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
add_rect(s, 0, 0, 13.33, 0.8, fill=DARK)
add_text(s, "СЕКЦІЇ РЕЗУЛЬТАТУ", 0.5, 0.2, 8, 0.5, size=11, bold=True, color=WHITE)
add_text(s, "Що є в боковій панелі", 0.5, 1.05, 12, 0.55, size=26, bold=True, color=DARK)

secs = [
    ("⚠",  ORANGE, "SOW + Delivery Spec",  "Таблиця параметрів здачі (роздільність,\nDPI, формат, час доби). Статус кожного:\nз брифу / дефолт / уточнити."),
    ("⚡",  RED,    "Конфлікти",            "Суперечності між файлами з обома\nджерелами і конкретним питанням\nдля уточнення."),
    ("▶",  BLUE,   "Роадмап",              "Впорядкований план роботи по стадіях\nвиробництва з задачами і нотатками\nпро залежності."),
    ("✓",  GREEN,  "Чеклист здачі",        "Плоский список всіх вимог по стадіях\nдля звірки результату перед здачею."),
    ("📋", DARK,   "Джерела",              "Посторінковий журнал — що знайдено\nв кожному файлі. Клік відкриває\nоригінал на потрібній сторінці."),
    ("💬", GRAY,   "Коментарі клієнта",    "Всі нотатки і підписи з файлів,\nзібрані дослівно з прив'язкою\nдо сторінки."),
]
for i, (ico, clr, title, desc) in enumerate(secs):
    col, row = i % 2, i // 2
    x = 0.4 + col * 6.5
    y = 1.85 + row * 1.8
    add_rect(s, x, y, 6.1, 1.6, fill=LIGHT, line=RGBColor(0xe0,0xe0,0xe0))
    add_rect(s, x, y, 0.45, 1.6, fill=clr)
    add_text(s, ico, x+0.07, y+0.5, 0.35, 0.5, size=16, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, title, x+0.6, y+0.15, 5.3, 0.4, size=13, bold=True, color=DARK)
    add_text(s, desc, x+0.6, y+0.6, 5.3, 0.9, size=10, color=GRAY)

# ─── SLIDE 7 — Delivery Spec ──────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
add_rect(s, 0, 0, 13.33, 0.8, fill=DARK)
add_text(s, "DELIVERY SPEC", 0.5, 0.2, 8, 0.5, size=11, bold=True, color=WHITE)
add_text(s, "Що буде здано клієнту — завжди зрозуміло", 0.5, 1.05, 12, 0.55, size=24, bold=True, color=DARK)
add_text(s, "Навіть якщо клієнт не вказав частину параметрів — дефолти підставляються автоматично з виробничих шаблонів.",
         0.5, 1.75, 12, 0.45, size=13, color=GRAY)

rows = [
    ("Роздільність", "4K",     "brief"),
    ("DPI",          "72 dpi", "default"),
    ("Формат файлу", "JPEG",   "default"),
    ("Час доби",     "вечір",  "brief"),
    ("Кількість зображень", "—", "unclear"),
    ("Співвідношення сторін", "16x9", "default"),
]
add_rect(s, 0.5, 2.4, 12.3, 0.38, fill=RGBColor(0xf0,0xee,0xea))
for j, h in enumerate(["Параметр", "Значення", "Статус"]):
    add_text(s, h, [0.65, 5.3, 9.5][j], 2.47, 3.5, 0.28, size=9, bold=True, color=GRAY)

for i, (param, val, src) in enumerate(rows):
    y = 2.85 + i * 0.52
    bg = LIGHT if i % 2 == 0 else RGBColor(0xff,0xff,0xff)
    add_rect(s, 0.5, y, 12.3, 0.5, fill=bg)
    add_text(s, param, 0.65, y+0.1, 4.4, 0.35, size=12, color=DARK)
    add_text(s, val,   5.3,  y+0.1, 3.8, 0.35, size=12, color=DARK if val != "—" else GRAY)
    if src == "brief":
        add_text(s, "✓ з брифу", 9.5, y+0.1, 2.5, 0.35, size=10, bold=True, color=GREEN)
    elif src == "default":
        add_text(s, "дефолт",    9.5, y+0.1, 2.5, 0.35, size=10, color=GRAY)
    else:
        add_text(s, "⚠ уточнити", 9.5, y+0.1, 2.5, 0.35, size=10, bold=True, color=ORANGE)

# ─── SLIDE 8 — SOW Matrix ─────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
add_rect(s, 0, 0, 13.33, 0.8, fill=DARK)
add_text(s, "SOW МАТРИЦЯ — ВКЛАДКА ЗВІТ", 0.5, 0.2, 8, 0.5, size=11, bold=True, color=WHITE)
add_text(s, "Покриття кожного пункту шаблону — одним поглядом", 0.5, 1.05, 12, 0.55, size=24, bold=True, color=DARK)
add_text(s, "Будується автоматично після основного аналізу (окремий легкий запит без файлів).",
         0.5, 1.72, 12, 0.4, size=13, color=GRAY)

sow_rows = [
    ("Креслення плану",    "found",   "DWG з розміщенням меблів",  "БРИФ стор.1"),
    ("Меблі та декор",     "found",   "Диван Minotti, стіл West Elm", "МАТЕРІАЛИ стор.3"),
    ("Оздоблення стін",    "partial", "Є колір, немає матеріалу",  "БРИФ стор.2"),
    ("Вид з вікна",        "missing", "—",                         "—"),
    ("Роздільність",       "found",   "4K вказано в брифі",        "БРИФ стор.1"),
    ("Наявність людей",    "missing", "—",                         "—"),
]
add_rect(s, 0.5, 2.3, 12.3, 0.38, fill=RGBColor(0xf0,0xee,0xea))
for j, h in enumerate(["Пункт SOW", "Статус", "Знайдено", "Джерело"]):
    add_text(s, h, [0.65, 4.5, 6.0, 10.5][j], 2.37, 3.5, 0.28, size=9, bold=True, color=GRAY)

for i, (item, status, found, source) in enumerate(sow_rows):
    y = 2.75 + i * 0.52
    bg = LIGHT if i % 2 == 0 else RGBColor(0xff,0xff,0xff)
    add_rect(s, 0.5, y, 12.3, 0.5, fill=bg)
    ico = "✅" if status == "found" else "⚠️" if status == "partial" else "❌"
    clr = GREEN if status == "found" else ORANGE if status == "partial" else RED
    add_text(s, item,   0.65, y+0.1, 3.6, 0.35, size=11, color=DARK)
    add_text(s, ico,    4.55, y+0.1, 1.2, 0.35, size=11, bold=True, color=clr)
    add_text(s, found,  6.0,  y+0.1, 4.2, 0.35, size=11, color=DARK if found != "—" else GRAY)
    sc = BLUE if source != "—" else GRAY
    add_text(s, source + (" ↗" if source != "—" else ""), 10.5, y+0.1, 2.1, 0.35, size=10, color=sc)

add_text(s, "3 знайдено  ·  1 неповно  ·  2 відсутнє", 0.65, 6.95, 6, 0.35, size=10, color=GRAY)

# ─── SLIDE 9 — Export ─────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
add_rect(s, 0, 0, 13.33, 0.8, fill=DARK)
add_text(s, "ЕКСПОРТ", 0.5, 0.2, 6, 0.5, size=11, bold=True, color=WHITE)
add_text(s, "Два формати вивантаження", 0.5, 1.1, 12, 0.6, size=28, bold=True, color=DARK)

add_rect(s, 0.5, 2.1, 5.9, 4.5, fill=LIGHT, line=RGBColor(0xe0,0xe0,0xe0))
add_rect(s, 0.5, 2.1, 5.9, 0.5, fill=GREEN)
add_text(s, "📊  Excel (XLS)", 0.75, 2.18, 5.4, 0.35, size=14, bold=True, color=WHITE)
add_text(s,
    "Кнопка у тулбарі режиму «Таблиця».\n\n"
    "Скачує tz-YYYY-MM-DD.xlsx з колонками:\n"
    "Тип · Вимога · Цитата · Категорія ·\n"
    "Кімната · Стадія · Джерело · Посилання\n\n"
    "Посилання клікабельні в Excel.\n"
    "Враховує активні фільтри.",
    0.75, 2.75, 5.4, 3.7, size=12, color=DARK)

add_rect(s, 6.9, 2.1, 5.9, 4.5, fill=LIGHT, line=RGBColor(0xe0,0xe0,0xe0))
add_rect(s, 6.9, 2.1, 5.9, 0.5, fill=BLUE)
add_text(s, "🖨  PDF (друк)", 7.15, 2.18, 5.4, 0.35, size=14, bold=True, color=WHITE)
add_text(s,
    "Кнопка «PDF» у хедері.\n\n"
    "window.print() з оптимізованими\nстилями друку:\n\n"
    "— Таблиця без розривів рядків\n"
    "— UI-елементи приховані\n"
    "— Посилання виводяться як текст\n"
    "— Назва файлу = тип проекту + дата",
    7.15, 2.75, 5.4, 3.7, size=12, color=DARK)

# ─── SLIDE 10 — Project types ─────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
add_rect(s, 0, 0, 13.33, 0.8, fill=DARK)
add_text(s, "ТИПИ ПРОЕКТІВ", 0.5, 0.2, 8, 0.5, size=11, bold=True, color=WHITE)
add_text(s, "9 типів — кожен зі своїм SOW-шаблоном і дефолтами", 0.5, 1.05, 12, 0.55, size=22, bold=True, color=DARK)

types = [
    ("🛋", "Інтер'єр житловий",       "4K · 72dpi · JPEG · день"),
    ("🏢", "Інтер'єр комерційний",     "4K · 72dpi · JPEG · день"),
    ("🏠", "Екстер'єр / фасад",        "4K · 72dpi · JPEG · літо/день\n+ Aerial підтип"),
    ("✨", "Лайфстайл рендеринг",      "4K · 72dpi · JPEG · sunset"),
    ("📦", "Silo рендеринг",           "4K · 300dpi · JPEG · 1x1"),
    ("🗺", "Мастерплан",               "—"),
    ("🪑", "Продуктова візуалізація",  "4K · 72dpi · JPEG"),
    ("🧊", "3D Моделювання",           ".max · Real World UV · Corona\n+ AR підтип"),
    ("📐", "Флорплан",                 "FullHD · 72dpi · PNG · top-down"),
]
for i, (ico, name, defaults) in enumerate(types):
    col, row = i % 3, i // 3
    x = 0.4 + col * 4.3
    y = 1.9 + row * 1.75
    add_rect(s, x, y, 4.0, 1.55, fill=LIGHT, line=RGBColor(0xe0,0xe0,0xe0))
    add_text(s, ico,      x+0.15, y+0.15, 0.6,  0.6,  size=22)
    add_text(s, name,     x+0.8,  y+0.15, 3.0,  0.45, size=12, bold=True, color=DARK)
    add_text(s, defaults, x+0.8,  y+0.65, 3.1,  0.75, size=10, color=GRAY)

# ─── SLIDE 11 — Tips ──────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
add_rect(s, 0, 0, 13.33, 0.8, fill=DARK)
add_text(s, "ПОРАДИ", 0.5, 0.2, 6, 0.5, size=11, bold=True, color=WHITE)
add_text(s, "Для кращого результату", 0.5, 1.05, 12, 0.55, size=28, bold=True, color=DARK)

tips = [
    ("📁", "Завантажуй всі файли одразу",    "Бриф + креслення + референси разом дають кращий результат.\nClause зіставляє їх між собою — чим більше контексту, тим точніше."),
    ("📐", "DWG краще за PDF-скан",           "З DWG витягуються назви приміщень, шари і розміри.\nPDF-скан дає тільки зображення без метаданих."),
    ("⚠",  "Перевіряй конфлікти першими",     "Суперечності між файлами краще уточнити до початку роботи.\nПізніше переробка коштує значно дорожче."),
    ("🗜",  "ZIP для великих пакетів",         "Архів розпаковується автоматично і класифікується.\nЗручно якщо файлів багато — одна дія замість десяти."),
    ("💾", "Сесія зберігається автоматично",  "Закрив вкладку — не страшно. Остання сесія відновлюється\nкнопкою на стартовому екрані."),
]
for i, (ico, title, desc) in enumerate(tips):
    y = 1.9 + i * 1.05
    add_rect(s, 0.5, y, 12.3, 0.95, fill=LIGHT, line=RGBColor(0xe8,0xe8,0xe8))
    add_text(s, ico,   0.7,  y+0.22, 0.5, 0.5, size=20)
    add_text(s, title, 1.4,  y+0.1,  4.5, 0.38, size=13, bold=True, color=DARK)
    add_text(s, desc,  1.4,  y+0.52, 10.9, 0.4, size=11, color=GRAY)

# ─── SLIDE 12 — End ───────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
add_rect(s, 0, 0, 13.33, 7.5, fill=DARK)
add_rect(s, 0, 0, 0.06, 7.5, fill=BLUE)
add_text(s, "Спробуй зараз", 0.6, 2.2, 12, 1.0, size=48, bold=True, color=WHITE)
add_text(s, "dlytvynov-rgb.github.io/tz-tool", 0.6, 3.5, 12, 0.6, size=22, color=BLUE)
add_text(s, "Web · macOS · Windows · Безкоштовно", 0.6, 4.3, 12, 0.45, size=14, color=GRAY)

out = "C:/Users/dima/OneDrive/Desktop/TZ_Tool_Presentation.pptx"
prs.save(out)
print(f"Saved: {out}")
